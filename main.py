import os
import random
import sys

from pptx import Presentation
from pptx.shapes.placeholder import PicturePlaceholder
from pptx.slide import Slide


# 创建并填充列表


def init_list(maximum=10) -> list:
    __itemList = []
    for __i in range(0, maximum):
        __itemList.append(__i)
    return __itemList


# 从列表中随机取出1个元素并从列表中删除
def get_random_item(itemList: list) -> any:
    __index = random.randint(0, len(itemList) - 1)
    __item = itemList[__index]
    del itemList[__index]
    return __item


# 打乱卡牌列表
def random_list(itemList: list) -> list:
    __tempItemList = itemList.copy()
    __newItemList = []
    while len(__tempItemList) > 0:
        __newItemList.append(get_random_item(__tempItemList))
    return __newItemList


# 以每局卡牌不重复分组
def group_cardList(itemList: list, maximum=10) -> list:
    __tempItemList = itemList.copy()
    __groupList = []
    while len(__tempItemList) >= maximum:
        __tempIndexList = random_list(init_list(maximum))
        __tempCardList = __tempItemList[:maximum]
        __groupList.append({
            "cardList": __tempCardList,  # 本轮的卡牌
            "murderer": __tempCardList[get_random_item(__tempIndexList)],  # 杀手鹅
            "victim": __tempCardList[get_random_item(__tempIndexList)]  # 受害鹅
        })
        del __tempItemList[:maximum]
    return __groupList


# 以每局杀手鹅和受害鹅不重复分组
def group_cardList2(itemList: list, maximum=10) -> list:
    __tempItemList = itemList.copy()
    __groupList = []
    while len(__tempItemList) >= maximum:
        __tempItemList = random_list(__tempItemList)
        __tempCardList = __tempItemList[:maximum]
        __murderer = __tempCardList[0]
        __victim = __tempCardList[1]
        __tempCardList = random_list(__tempCardList)

        __groupList.append({
            "cardList": __tempCardList,  # 本轮的卡牌
            "murderer": __murderer,  # 杀手鹅
            "victim": __victim  # 受害鹅
        })
        del __tempItemList[:2]
    return __groupList


# 按给定的母版编号添加页面
def pptx_addSlide(__pptx: Presentation, layoutIndex: int) -> Slide:
    return __pptx.slides.add_slide(__pptx.slide_layouts[layoutIndex])


def pptx_fillWitness(__slide: Slide, group: dict):
    __murderer = group["murderer"]
    __victim = group["victim"]
    __slide.placeholders[10].insert_picture(f"./cards/{__murderer}")
    __slide.placeholders[11].insert_picture(f"./cards/{__victim}")


# 填充指认界面的图像
def pptx_fillIdentify(__slide: Slide, group: dict):
    __cardList = group["cardList"]
    __murderer = group["murderer"]
    __victim = group["victim"]

    for __i in range(10):
        __shape: PicturePlaceholder = __slide.placeholders[__i + 10]
        __shape.insert_picture(f"./cards/{__cardList[__i]}")
        if __cardList[__i] == __murderer:
            __slide.placeholders[__i + 20].insert_picture("./images/murderer.png")
        if __cardList[__i] == __victim:
            __slide.placeholders[__i + 20].insert_picture("./images/victim.png")


if __name__ == '__main__':
    mode = "pptx"  # 生成模式 pptx 或 mp4
    group_rule = "all"  # 分组规则 all:每局卡牌全部不重复,role:身份牌不重复

    if len(sys.argv) > 1:
        try:
            mode = sys.argv[1].lower()
            group_rule = sys.argv[2].lower()
        except Exception:
            pass
        print(f"模式:{mode}, 规则:{group_rule}")

    cardList = random_list(os.listdir("./cards"))
    roundList = group_cardList(cardList)
    if group_rule == "role":
        roundList = group_cardList2(cardList)

    try:
        os.mkdir("./pptxs")
        os.mkdir("./videos")
    except FileExistsError:
        pass

    start_round = int(input("请输入开始时的局号:"))
    if mode == "pptx":
        max_round = int(input("请输入PPTX内的最大局数:"))
        pptx: Presentation = Presentation("./template/以鹅传鹅.pptx")
        i = 0
        for i in range(len(roundList)):
            if i % max_round == 0:
                start_round = i + 1
                pptx: Presentation = Presentation("./template/以鹅传鹅.pptx")
                pptx_addSlide(pptx, 0)

            print(f"添加 第{i % max_round + start_round:02d}局 游戏")
            pptx_fillWitness(pptx_addSlide(pptx, 1), roundList[i])
            pptx_addSlide(pptx, 2)
            pptx_fillIdentify(pptx_addSlide(pptx, 3), roundList[i])

            if i % max_round == max_round - 1:
                pptx_addSlide(pptx, 4)
                print("保存 演示文稿\r\n")
                pptx.save(f"./pptxs/第{start_round:02d}-{i % max_round + start_round:02d}局.pptx")
            elif i == len(roundList) - 1:
                pptx_addSlide(pptx, 4)
                print("保存 演示文稿\r\n")
                pptx.save(f"./pptxs/第{start_round:02d}-{i % max_round + start_round:02d}局.pptx")
    elif mode == "mp4":
        # 生成视频用PPT
        max_round = int(input("请输入游戏总局数:"))
        for i in range(len(roundList)):
            if max_round <= 0:
                max_round = len(roundList)
            if i < max_round:
                print(f"输出 第{i % max_round + start_round:02d}局 游戏")
                pptx_video: Presentation = Presentation("./template/以鹅传鹅.pptx")
                pptx_fillWitness(pptx_addSlide(pptx_video, 5), roundList[i])
                pptx_video.save(f"./videos/第{start_round + i:02d}局_01目击.pptx")

                pptx_video: Presentation = Presentation("./template/以鹅传鹅.pptx")
                pptx_fillIdentify(pptx_addSlide(pptx_video, 6), roundList[i])
                pptx_video.save(f"./videos/第{start_round + i:02d}局_03指认.pptx")

                pptx_video: Presentation = Presentation("./template/以鹅传鹅.pptx")
                pptx_fillIdentify(pptx_addSlide(pptx_video, 7), roundList[i])
                pptx_video.save(f"./videos/第{start_round + i:02d}局_04结果.pptx")
